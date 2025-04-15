from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

def create_title_slide(prs, title, subtitle):
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = subtitle

def create_content_slide(prs, title, content):
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    content_slide.shapes.title.text = title
    content_slide.placeholders[1].text = content

def main():
    # 创建演示文稿
    prs = Presentation()
    
    # 设置标题页
    create_title_slide(prs, 
                      "基于AutoGen的多智能体代码转换系统",
                      "一个智能的代码语言转换框架")
    
    # 目录页
    create_content_slide(prs, "目录", 
                        "1. 系统概述\n"
                        "2. 核心组件\n"
                        "3. 工作流程\n"
                        "4. 技术特点\n"
                        "5. 使用示例")
    
    # 系统概述
    create_content_slide(prs, "1. 系统概述",
                        "• 基于AutoGen框架\n"
                        "• 多智能体协作系统\n"
                        "• 支持多语言代码转换\n"
                        "• 本地模型部署")
    
    # 核心组件
    create_content_slide(prs, "2. 核心组件",
                        "专家智能体：\n"
                        "• CodeAnalyzerAgent：代码分析专家\n"
                        "• ImportConverterAgent：导入转换专家\n"
                        "• ClassConverterAgent：类转换专家\n"
                        "• FunctionConverterAgent：函数转换专家\n"
                        "• MainConverterAgent：主程序转换专家\n"
                        "• CodeIntegratorAgent：代码整合专家\n\n"
                        "系统组件：\n"
                        "• UserProxyAgent：用户代理\n"
                        "• GroupChat：群聊系统\n"
                        "• GroupChatManager：群聊管理器")
    
    # 工作流程
    create_content_slide(prs, "3. 工作流程",
                        "1. 输入源代码和目标语言\n"
                        "2. 加载转换规则\n"
                        "3. 多智能体协作转换\n"
                        "   • 分析代码结构\n"
                        "   • 转换各个组件\n"
                        "   • 整合最终代码\n"
                        "4. 输出转换结果")
    
    # 技术特点
    create_content_slide(prs, "4. 技术特点",
                        "• 模块化设计\n"
                        "• 多语言支持\n"
                        "• 可扩展规则系统\n"
                        "• 本地模型部署\n"
                        "• 多智能体协作")
    
    # 使用示例
    create_content_slide(prs, "5. 使用示例",
                        "```python\n"
                        "transformer = AutoGenTransformer()\n"
                        "converted_code = transformer.convert_code(\n"
                        "    code=source_code,\n"
                        "    source_lang=\"Python\",\n"
                        "    target_lang=\"Java\"\n"
                        ")\n"
                        "```")
    
    # 总结
    create_content_slide(prs, "总结",
                        "• 高效的多智能体协作系统\n"
                        "• 灵活的代码转换框架\n"
                        "• 本地化部署方案\n"
                        "• 可扩展的架构设计")
    
    # 获取当前脚本所在目录
    current_dir = os.path.dirname(os.path.abspath(__file__))
    # 构建PPT文件的完整路径
    ppt_path = os.path.join(current_dir, 'code_transformer.pptx')
    
    # 保存PPT
    prs.save(ppt_path)
    print(f"PPT文件已生成在: {ppt_path}")

if __name__ == '__main__':
    main() 